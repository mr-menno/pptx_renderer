import pytest
import json
from pathlib import Path
from pptx import Presentation
from tempfile import TemporaryDirectory

@pytest.fixture
def input_json():
    data = {
        "myvar1": "hello 1",
        "myvar2": "hello 2",
    }
    input_path = Path("input.json")
    input_path.write_text(json.dumps(data))
    yield input_path
    input_path.unlink()

@pytest.fixture
def input_template():
    """Input template for testing"""
    temp_dir = TemporaryDirectory()
    # Create a presentation object from temp_path and add a slide
    # in the slide, add two textboxes with the text "{{{myvar1}}}" and "{{{myvar2}}}"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = "{{{myvar1}}}"
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = "{{{myvar2}}}"
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir) / "template.pptx"
        prs.save(temp_path)
        yield temp_path