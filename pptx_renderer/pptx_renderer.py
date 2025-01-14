"""Main Module"""

import re
from os import PathLike
from pathlib import Path
from typing import Any, Dict, Optional, Union
from warnings import warn as warning
from functools import partial
from . import plugins

from pptx import Presentation

from .exceptions import RenderError
from .utils import fix_quotes, para_text_replace

PLUGINS = [plugins.image, plugins.video, plugins.table]

class PPTXRenderer:
    """PPTX Renderer class

    This class is used to render a PPTX template by replacing python statements
    with the result of evaluating the python statements.

    Attributes:
        template_path (str): Path to the PPTX template.
    """

    def __init__(self, template_path: Union[str, PathLike]):
        self.template_path = template_path
        self.plugins = {}
        self.namespace = {}
        for plugin in PLUGINS:
            self.register_plugin(plugin.__name__, plugin)
    
    def register_plugin(self, name: str, func: callable):
        """Register a plugin function.

        The plugin function should take 2 or more arguments. The first argument
        is the result of evaluating the python statement. The second argument is
        a dictionary containing the following keys:
        - shape: The pptx shape object where the placeholder is present
        - slide: The pptx slide object where the placeholder is present
        - slide_no: The slide number where the placeholder is present
        The remaining arguments are the arguments passed to the plugin function

        Args:
            name (str): Name of the plugin.
            func (callable): Function to be registered.

        Returns:
            None
        """
        self.plugins[name] = func

    def render(
        self,
        output_path: Union[str, PathLike],
        methods_and_params: Optional[Dict[str, Any]] = None,
        skip_failed: bool = False,
    ) -> None:
        """Render PPTXRenderer template and save to output_path.

        Args:
            output_path (str): Path to the output PPTX file.
            methods_and_params (dict, optional): Dictionary of methods and parameters
                to be used in the template. Defaults to None.
            skip_failed (bool, optional): Dont raise an error if some of the
                statements failed to render. Defaults to False.

        Returns:
            None
        """
        if not Path(self.template_path).exists():
            raise (FileNotFoundError(f"{self.template_path} not found"))
        outppt = Presentation(self.template_path)
        self.namespace.update(methods_and_params)
        for slide_no, slide in enumerate(outppt.slides):
            if slide.has_notes_slide:
                python_code = re.search(
                    r"```python([\s\S]*)```",
                    fix_quotes(slide.notes_slide.notes_text_frame.text),
                    re.MULTILINE,
                )
                if python_code:
                    exec(python_code.group(1), self.namespace)
            for shape in list(slide.shapes):
                if shape.has_text_frame:
                    matches = re.finditer(r"{{{(.*)}}}", shape.text)
                    if not matches:
                        continue
                    for match_assignment in matches:
                        parts = match_assignment.group(1).split(":")
                        try:
                            result = eval(fix_quotes(parts[0]), self.namespace)
                        except Exception as ex:
                            if skip_failed:
                                warning(
                                    f"Evaluation of '{parts[0]}' in slide {slide_no+1} failed"
                                )
                                continue
                            raise RenderError(
                                f"Failed to evaluate '{parts[0]}'."
                            ) from ex
                        if len(parts) > 1:
                            namespace = self.namespace.copy()
                            context = {
                                "result": result,
                                "presentation": outppt,
                                "shape": shape,
                                "slide": slide,
                                "slide_no": slide_no,
                            }
                            for plugin_name, plugin in self.plugins.items():
                                func = partial(plugin, context)
                                namespace[plugin_name] = func 
                            try:
                                exec(fix_quotes(parts[1]), namespace)
                            except Exception as ex:
                                if skip_failed:
                                    warning(
                                        f"Failed to render {parts[0]} in slide {slide_no+1}"
                                    )
                                    continue
                                raise RenderError(
                                    f"Failed to render {parts[0]} in slide {slide_no+1}"
                                ) from ex
                        else:
                            for paragraph in shape.text_frame.paragraphs:
                                para_text_replace(
                                    paragraph, match_assignment.group(0), result
                                )
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            matches = re.finditer(r"{{{(.*)}}}", cell.text)
                            if not matches:
                                continue
                            for match_assignment in matches:
                                parts = match_assignment.group(1).split(":")
                                try:
                                    result = eval(fix_quotes(parts[0]), self.namespace)
                                except Exception as ex:
                                    if skip_failed:
                                        warning(
                                            f"Evaluation of '{parts[0]}' in slide {slide_no+1} failed"
                                        )
                                        continue
                                    raise RenderError(
                                        f"Failed to evaluate '{parts[0]}'."
                                    ) from ex
                                for paragraph in cell.text_frame.paragraphs:
                                    para_text_replace(
                                        paragraph, match_assignment.group(0), result
                                    )
        outppt.save(output_path)
